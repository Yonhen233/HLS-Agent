from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "docs" / "presentations"
PPTX_PATH = OUT_DIR / "HLS_Agent_MNIST_Project_Pitch_20260617.pptx"
NOTES_PATH = OUT_DIR / "HLS_Agent_MNIST_Project_Pitch_20260617_notes.md"


COLORS = {
    "navy": RGBColor(20, 40, 67),
    "blue": RGBColor(40, 94, 145),
    "teal": RGBColor(18, 151, 147),
    "green": RGBColor(67, 160, 71),
    "orange": RGBColor(239, 136, 53),
    "red": RGBColor(201, 73, 62),
    "ink": RGBColor(38, 43, 51),
    "muted": RGBColor(96, 105, 118),
    "light": RGBColor(242, 246, 250),
    "white": RGBColor(255, 255, 255),
    "line": RGBColor(207, 216, 225),
}


FONT = "Microsoft YaHei"


def set_text(frame, text: str, size: int = 18, bold: bool = False, color=None, align=None) -> None:
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color or COLORS["ink"]
    if align:
        p.alignment = align


def add_textbox(slide, x, y, w, h, text, size=18, bold=False, color=None, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(box.text_frame, text, size=size, bold=bold, color=color, align=align)
    return box


def add_title(slide, title: str, subtitle: str | None = None):
    add_textbox(slide, 0.65, 0.32, 11.8, 0.45, title, size=27, bold=True, color=COLORS["navy"])
    underline = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.65),
        Inches(0.92),
        Inches(1.2),
        Inches(0.05),
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = COLORS["teal"]
    underline.line.fill.background()
    if subtitle:
        add_textbox(slide, 0.65, 1.02, 11.6, 0.32, subtitle, size=11, color=COLORS["muted"])


def add_footer(slide, page: int):
    add_textbox(slide, 0.65, 7.18, 8.5, 0.25, "DL-Operator-to-HLS Agent · MNIST 真实识别 Demo · 2026-06-17", size=8, color=COLORS["muted"])
    add_textbox(slide, 12.0, 7.18, 0.55, 0.25, str(page), size=8, color=COLORS["muted"], align=PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, title, body, accent="teal"):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["white"]
    shape.line.color.rgb = COLORS["line"]
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS[accent]
    bar.line.fill.background()
    add_textbox(slide, x + 0.18, y + 0.13, w - 0.32, 0.32, title, size=15, bold=True, color=COLORS["navy"])
    add_textbox(slide, x + 0.18, y + 0.52, w - 0.32, h - 0.6, body, size=11, color=COLORS["ink"])


def add_table(slide, x, y, w, h, rows, cols, data, header_fill="navy", font_size=9):
    table_shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT
                p.font.size = Pt(font_size if r else font_size + 1)
                p.font.bold = r == 0
                p.font.color.rgb = COLORS["white"] if r == 0 else COLORS["ink"]
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS[header_fill] if r == 0 else (RGBColor(248, 251, 253) if r % 2 == 0 else COLORS["white"])
    return table_shape


def add_flow_box(slide, x, y, w, h, text, color="blue"):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[color]
    shape.line.fill.background()
    set_text(shape.text_frame, text, size=10, bold=True, color=COLORS["white"], align=PP_ALIGN.CENTER)
    return shape


def add_arrow(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = COLORS["muted"]
    line.line.width = Pt(1.6)
    return line


def new_slide(prs, title: str, subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(250, 252, 253)
    add_title(slide, title, subtitle)
    return slide


def read_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def build_deck() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    report = read_json(ROOT / "runs" / "mnist_recognition_mlp_234d539d" / "report.json")
    verification = read_json(ROOT / "runs" / "mnist_recognition_mlp_234d539d" / "verification.json")
    training = read_json(ROOT / "models" / "mnist_recognition" / "mnist_mlp_training_metrics.json")
    cls = verification["classification"]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    notes: list[str] = []

    # 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["navy"]
    add_textbox(slide, 0.85, 0.65, 11.5, 0.72, "深度学习算子转 HLS Agent", size=34, bold=True, color=COLORS["white"])
    add_textbox(slide, 0.9, 1.45, 11.4, 0.35, "面向 FPGA HLS 工作流的 Tool-using Agent 原型", size=16, color=RGBColor(214, 231, 242))
    add_textbox(slide, 0.9, 2.35, 10.8, 0.6, "用 MNIST 真实识别 Demo 证明：模型 → HLS → Vivado csim/csynth → 识别验证，已经形成可追踪闭环。", size=22, bold=True, color=COLORS["white"])
    add_card(slide, 0.9, 4.05, 3.7, 1.3, "真实结果", "HLS csim accuracy 95%\nargmax match 100%\ntiming met", "green")
    add_card(slide, 4.85, 4.05, 3.7, 1.3, "工程框架", "Todo + Specialist\nToolRegistry + Trace\nMemory/RAG", "teal")
    add_card(slide, 8.8, 4.05, 3.7, 1.3, "下一阶段", "CNN adapter\nprecision advisor\n上板前 IP 化", "orange")
    add_textbox(slide, 0.9, 6.85, 10.5, 0.25, "汇报目标：争取导师支持继续推进“模型到 HLS”的工程化 Agent 方向", size=10, color=RGBColor(214, 231, 242))
    notes.append("开场强调：这不是单纯写脚本，而是一个可验证、可追踪、能沉淀经验的 HLS Agent 原型。")

    # 2
    slide = new_slide(prs, "为什么这个项目值得做", "痛点：深度学习模型到 FPGA HLS 的链路长、工具碎、失败难定位")
    add_card(slide, 0.75, 1.55, 3.8, 1.25, "工具链门槛高", "hls4ml / ONNX / Vivado HLS / report parser 分散，学生很难一次跑通。", "blue")
    add_card(slide, 4.8, 1.55, 3.8, 1.25, "失败路径复杂", "算子不支持、版本兼容、csim 数值误差、timing 失败都需要结构化处理。", "red")
    add_card(slide, 8.85, 1.55, 3.8, 1.25, "经验难复用", "precision、reuse_factor、clock、错误修复经验往往散落在日志里。", "orange")
    add_textbox(slide, 0.85, 3.3, 11.7, 0.4, "项目目标：把“模型到 HLS”的流程变成一个可追踪、可验证、可记忆、可扩展的 Agent 系统。", size=19, bold=True, color=COLORS["navy"])
    add_flow_box(slide, 1.0, 4.35, 1.65, 0.65, "用户任务", "blue")
    add_arrow(slide, 2.7, 4.68, 3.15, 4.68)
    add_flow_box(slide, 3.2, 4.35, 1.85, 0.65, "Agent 编排", "teal")
    add_arrow(slide, 5.1, 4.68, 5.55, 4.68)
    add_flow_box(slide, 5.6, 4.35, 2.0, 0.65, "EDA 工具链", "orange")
    add_arrow(slide, 7.65, 4.68, 8.1, 4.68)
    add_flow_box(slide, 8.15, 4.35, 1.85, 0.65, "报告解析", "green")
    add_arrow(slide, 10.05, 4.68, 10.5, 4.68)
    add_flow_box(slide, 10.55, 4.35, 1.85, 0.65, "经验沉淀", "blue")
    add_textbox(slide, 1.0, 5.55, 11.1, 0.5, "价值：既服务 FPGA/HLS 研究，也能体现现代 Agent 工程能力。", size=18, bold=True, color=COLORS["teal"])
    add_footer(slide, 2)
    notes.append("这一页向导师说明研究动机：工具链复杂和经验难沉淀，是 Agent 可以发挥价值的地方。")

    # 3
    slide = new_slide(prs, "项目基本原理", "Main Agent 负责编排，Specialist 处理局部领域任务，工具负责事实")
    boxes = [
        ("Main Agent", "任务理解\n路径选择\n状态合并\n最终报告", "navy"),
        ("Todo Runtime", "Plan → Todo\nExecute → Reflect\nFinalize", "blue"),
        ("Specialists", "HLS4ML\nVivado\nVerification\nOptimization\nMemory", "teal"),
        ("Tools", "hls4ml\nVivado HLS\nReport parser\nDB/RAG", "orange"),
        ("Artifacts/Memory", "HLS 工程\nlogs/reports\nSQLite\nRAG chunks", "green"),
    ]
    x = 0.85
    for title, body, color in boxes:
        add_card(slide, x, 1.65, 2.25, 2.1, title, body, color)
        x += 2.48
    add_textbox(slide, 1.0, 4.35, 11.3, 0.55, "核心思想：LLM handles ambiguity，Tools handle facts，Workflow handles guarantees，Verification decides trust。", size=18, bold=True, color=COLORS["navy"])
    add_table(slide, 1.0, 5.2, 11.2, 1.35, 4, 3, [
        ["层次", "负责什么", "为什么这样设计"],
        ["确定性 Workflow", "HLS 必经阶段", "保证可复现"],
        ["Agent 决策点", "路径选择/失败恢复/经验复用", "根据 observation 更新状态"],
        ["LLM 能力", "候选代码/repair/解释优化", "只在开放问题启用"],
    ], font_size=10)
    add_footer(slide, 3)
    notes.append("强调：没有把所有事情交给 LLM。确定性的地方工程化，开放性问题才交给 LLM。")

    # 4
    slide = new_slide(prs, "当前工作流", "以 MNIST 真实识别 Demo 为例")
    steps = [
        ("训练 MLP", "PyTorch\n91.76% eval"),
        ("导出 ONNX", "Gemm + Relu\n静态 shape"),
        ("hls4ml adapter", "Gemm → Dense\nRelu → Activation"),
        ("生成 HLS 工程", "hls4ml\nVivado backend"),
        ("Vivado HLS", "csim + csynth\n2018.3"),
        ("验证/报告", "accuracy\nlatency/resource"),
        ("Memory", "verified impl\nparameter exp"),
    ]
    x = 0.55
    for i, (title, body) in enumerate(steps):
        add_flow_box(slide, x, 2.0, 1.55, 0.8, f"{title}\n{body}", "teal" if i not in {0, 4, 6} else ("blue" if i == 0 else "orange" if i == 4 else "green"))
        if i < len(steps) - 1:
            add_arrow(slide, x + 1.58, 2.4, x + 1.95, 2.4)
        x += 1.82
    add_table(slide, 1.0, 4.0, 11.2, 1.7, 4, 4, [
        ["阶段", "输入", "输出", "证据文件"],
        ["模型阶段", "MNIST 数据集", "ONNX + reference", "models/mnist_recognition/*"],
        ["HLS 阶段", "ONNX / task JSON", "HLS project", "runs/.../hls_project"],
        ["验证阶段", "csim/csynth", "accuracy + report", "summary.md / report.json"],
    ], font_size=9)
    add_textbox(slide, 1.0, 6.25, 11.3, 0.3, "本次不是 mock：hls4ml 与 Vivado HLS 2018.3 均走真实工具链。", size=16, bold=True, color=COLORS["green"])
    add_footer(slide, 4)
    notes.append("这里按实际执行路径讲，让导师能看到完整闭环。")

    # 5
    slide = new_slide(prs, "MNIST Demo：模型与验证集", "从随机结构 demo 升级为真实识别任务")
    add_table(slide, 0.9, 1.55, 6.1, 2.0, 5, 2, [
        ["项目", "结果"],
        ["模型结构", training["architecture"]],
        ["训练 epoch", len(training["history"])],
        ["5000 样本评估准确率", f"{training['best_eval_accuracy']*100:.2f}%"],
        ["HLS reference 样本数", training["reference_samples"]],
    ], font_size=11)
    add_table(slide, 7.35, 1.55, 5.0, 2.0, 5, 2, [
        ["文件", "作用"],
        ["mnist_mlp_trained.onnx", "HLS 转换输入"],
        ["mnist_test_inputs_20.dat", "csim 输入"],
        ["mnist_test_labels_20.json", "分类 label"],
        ["training_metrics.json", "训练证据"],
    ], font_size=9)
    add_card(slide, 0.95, 4.25, 5.45, 1.35, "为什么不直接用网上预训练 CNN", "外部 mnist-8.onnx 已下载，但包含 Conv/MaxPool/Reshape 等更多前端兼容风险；第一版主 demo 选择 HLS-friendly MLP，保证真实识别闭环稳定。", "orange")
    add_card(slide, 6.9, 4.25, 5.45, 1.35, "为什么仍然真实", "权重来自 MNIST 训练，不是随机初始化；输入来自真实 MNIST 测试集；HLS 输出与 ONNX reference 和 label 对比。", "green")
    add_footer(slide, 5)
    notes.append("回答导师可能的问题：为什么不用网上模型？重点是先保证可控、可复现的真实识别闭环。")

    # 6
    slide = new_slide(prs, "MNIST Demo：真实验证结果", "HLS 代码确实完成了数字识别")
    add_table(slide, 0.9, 1.45, 5.7, 2.3, 6, 2, [
        ["识别指标", "结果"],
        ["样本数", cls["sample_count"]],
        ["ONNX reference accuracy", f"{cls['reference_accuracy']*100:.0f}%"],
        ["HLS csim accuracy", f"{cls['hls_accuracy']*100:.0f}%"],
        ["Argmax match rate", f"{cls['argmax_match_rate']*100:.0f}%"],
        ["HLS correct", f"{cls['hls_correct']} / {cls['sample_count']}"],
    ], font_size=12)
    cmp = verification["comparison"]
    add_table(slide, 7.0, 1.45, 5.2, 2.3, 6, 2, [
        ["数值指标", "结果"],
        ["max abs error", cmp["max_abs_error"]],
        ["numeric tolerance", cmp["tolerance"]],
        ["numeric_passed", str(cmp["numeric_passed"])],
        ["recognition_passed", str(cmp["recognition_passed"])],
        ["verification status", verification["status"]],
    ], font_size=10)
    add_card(slide, 0.95, 4.45, 11.4, 1.25, "关键解释", "fixed-point HLS logits 与 float ONNX logits 数值漂移较大，但 10 类最大值位置保持一致，因此分类语义通过。报告同时保留 numeric drift 和 recognition pass，不掩盖误差。", "teal")
    add_footer(slide, 6)
    notes.append("这一页是最有说服力的数据页：HLS accuracy 95%，argmax match 100%。")

    # 7
    slide = new_slide(prs, "MNIST Demo：Vivado HLS 综合结果", "功能通过之外，也拿到了可解析的 HLS report")
    res = report["resources"]
    timing = report["timing"]
    add_table(slide, 0.9, 1.45, 11.4, 2.35, 7, 4, [
        ["类别", "指标", "结果", "说明"],
        ["Latency", "min / max", f"{report['latency']['min_cycles']} / {report['latency']['max_cycles']}", "cycles"],
        ["Interval", "min / max", f"{report['interval']['min_ii']} / {report['interval']['max_ii']}", "当前 Resource 策略"],
        ["Resources", "DSP / BRAM", f"{res['dsp']} / {res['bram']}", "乘法器/存储"],
        ["Resources", "LUT / FF", f"{res['lut']} / {res['ff']}", "逻辑资源"],
        ["Timing", "target / estimated", f"{timing['target_ns']}ns / {timing['estimated_ns']}ns", "满足约束"],
        ["Status", "timing met", str(timing["met"]), "可作为 HLS candidate"],
    ], font_size=9)
    add_card(slide, 0.95, 4.55, 3.5, 1.15, "结论 1", "真实 Vivado HLS 2018.3 完成 csim/csynth。", "green")
    add_card(slide, 4.9, 4.55, 3.5, 1.15, "结论 2", "report parser 成功提取 latency/resource/timing。", "teal")
    add_card(slide, 8.85, 4.55, 3.5, 1.15, "结论 3", "当前达到 deployment_ready_candidate。", "orange")
    add_footer(slide, 7)
    notes.append("强调这不是只生成代码，还完成综合并解析报告。")

    # 8
    slide = new_slide(prs, "已完成能力矩阵", "从 Demo 到 Agent 工程能力的映射")
    add_table(slide, 0.65, 1.35, 12.0, 4.95, 9, 4, [
        ["能力", "状态", "证据", "价值"],
        ["真实 hls4ml path", "已跑通", "MNIST MLP", "模型到 HLS"],
        ["Vivado HLS 2018.3", "已跑通", "csim/csynth", "真实 EDA toolchain"],
        ["Functional verification", "已跑通", "accuracy/argmax", "证明功能正确"],
        ["Report parser", "已跑通", "latency/resource/timing", "可量化优化"],
        ["Specialist isolation", "已接入", "HLS4ML/Vivado/Memory", "上下文隔离"],
        ["Memory/RAG", "已接入", "verified memory", "参数经验复用"],
        ["LLM candidate", "已验证", "Dense/MatMul/ReLU/Add", "开放算子生成"],
        ["Boundary unsupported", "已支持", "ResNet boundary", "不盲目承诺"],
    ], font_size=8)
    add_textbox(slide, 0.8, 6.55, 11.8, 0.35, "当前已经不是概念阶段：有真实模型、真实工具链、真实 report、真实 verification。", size=16, bold=True, color=COLORS["green"])
    add_footer(slide, 8)
    notes.append("这一页给导师总体看：不是只做了 MNIST，一套 Agent 工程能力已经成型。")

    # 9
    slide = new_slide(prs, "LLM Candidate 路径验证", "LLM 不是主路径必需品，而是 unsupported/candidate/repair 的增强能力")
    add_table(slide, 0.65, 1.35, 12.05, 3.1, 6, 7, [
        ["Demo", "Path", "Verification", "Latency", "DSP", "FF", "LUT"],
        ["Dense LLM", "llm_candidate", "golden csim passed", 37, 16, 1171, 2873],
        ["MatMul LLM", "llm_candidate", "golden csim passed", 3073, 8, 213, 520],
        ["ReLU LLM", "llm_candidate", "golden csim passed", 10, 0, 74, 300],
        ["Add LLM", "llm_candidate", "golden csim passed", 20, 0, 230, 205],
        ["ScaleShift LLM", "llm_candidate", "golden csim passed", 19, 0, 38, 111],
    ], font_size=8)
    add_card(slide, 0.85, 4.85, 5.45, 1.15, "设计原则", "确定性路径不强依赖 LLM；LLM 用于不确定路径、候选 HLS 生成、repair 和优化解释。", "blue")
    add_card(slide, 6.9, 4.85, 5.45, 1.15, "验证原则", "LLM 生成代码不能直接信任，必须通过 sandbox、testbench、csim、csynth、report parser。", "red")
    add_footer(slide, 9)
    notes.append("说明 LLM 的价值：不是每一步都调用，而是在开放问题上增强。")

    # 10
    slide = new_slide(prs, "哪些是写死 workflow，哪些体现 Agent 决策", "导师常问：这是不是普通 if-else 脚本？")
    add_table(slide, 0.75, 1.4, 5.65, 3.0, 6, 2, [
        ["写死的工程骨架", "为什么固定"],
        ["validate task", "输入安全"],
        ["inspect/check support", "HLS 必经阶段"],
        ["convert/synth/parse", "EDA 流程确定"],
        ["summary/memory", "可追踪沉淀"],
        ["tool permission", "安全边界"],
    ], font_size=10)
    add_table(slide, 6.9, 1.4, 5.65, 3.0, 6, 2, [
        ["Agent 决策点", "依据"],
        ["选择 hls4ml/fallback/LLM", "support result"],
        ["Gemm unsupported 后启用 adapter", "ONNX graph"],
        ["numeric fail 但 recognition pass", "verification metrics"],
        ["参数推荐", "verified memory"],
        ["memory promotion", "csim/csynth/timing"],
    ], font_size=10)
    add_card(slide, 0.95, 5.05, 11.4, 1.0, "推荐表述", "写死的是安全边界和工程流程；Agent 化的是状态驱动的路径选择、工具编排、失败恢复、经验复用和验证闭环。", "teal")
    add_footer(slide, 10)
    notes.append("这页直接回答导师/面试官的质疑：不是纯 if-else，也不是纯 LLM。")

    # 11
    slide = new_slide(prs, "项目当前进度", "有数据支撑的阶段性成果")
    add_table(slide, 0.6, 1.25, 12.1, 4.55, 8, 5, [
        ["方向", "当前状态", "代表 Demo", "关键数据", "成熟度"],
        ["Operator fallback", "已跑通", "Dense/MatMul", "真实 Vivado report", "稳定"],
        ["LLM candidate", "已跑通", "5 个算子", "均 golden csim passed", "可展示"],
        ["Model hls4ml", "已跑通", "MNIST MLP", "HLS acc 95%", "核心亮点"],
        ["Functional verification", "已增强", "classification-aware", "argmax 100%", "可复用"],
        ["Memory/RAG", "已接入", "ParameterAdvisor", "verified_history", "可扩展"],
        ["CNN/QONNX", "部分可用", "Tiny/QONNX", "前端仍需增强", "下一阶段"],
        ["Board deployment", "未开始", "待定板卡", "不生成 bitstream", "研究后续"],
    ], font_size=8)
    add_textbox(slide, 0.85, 6.25, 11.8, 0.4, "结论：当前已经完成“工程原型 + 真实 HLS 验证”的 MVP，下一阶段适合在导师支持下走向 CNN/上板前准备。", size=15, bold=True, color=COLORS["navy"])
    add_footer(slide, 11)
    notes.append("导师最关心进度：用成熟度表格说明哪些已经稳、哪些还在攻关。")

    # 12
    slide = new_slide(prs, "主要技术风险与应对", "不是回避问题，而是把风险工程化管理")
    add_table(slide, 0.8, 1.35, 11.8, 4.5, 6, 4, [
        ["风险", "当前现象", "影响", "应对策略"],
        ["工具链版本", "Vivado 2018.3 更稳，Vitis 新版不一定更好", "结果不可比", "固定基线，做隔离实验"],
        ["ONNX 前端边界", "CNN/Shape/Reshape 更复杂", "模型支持受限", "逐步扩展 adapter"],
        ["数值漂移", "fixed-point logits drift", "影响可信度", "区分 numeric / recognition"],
        ["资源压力", "MNIST MLP LUT/BRAM 不低", "板卡选择受限", "precision/reuse advisor"],
        ["上板未完成", "当前不生成 bitstream", "离部署有距离", "下一阶段 IP/AXI/board flow"],
    ], font_size=8)
    add_card(slide, 1.0, 6.1, 11.2, 0.72, "策略：先把模型到 HLS 的验证链路做扎实，再逐步进入 CNN、接口综合、IP packaging 和上板。", "orange")
    add_footer(slide, 12)
    notes.append("风险必须主动讲：导师更容易相信你知道问题在哪里。")

    # 13
    slide = new_slide(prs, "希望导师支持的方向", "需要资源与方向指导，而不是只要时间")
    add_card(slide, 0.85, 1.45, 3.65, 1.45, "研究方向确认", "是否聚焦“模型到 HLS Agent”作为主线；是否以 MNIST/CNN/算子库作为阶段目标。", "blue")
    add_card(slide, 4.85, 1.45, 3.65, 1.45, "板卡与工具链", "确定目标 FPGA 板卡；确认 Vivado/Vitis 稳定版本；建立共享验证环境。", "teal")
    add_card(slide, 8.85, 1.45, 3.65, 1.45, "评价标准", "定义 accuracy、latency、resource、timing、workflow success rate 等指标。", "green")
    add_card(slide, 0.85, 3.65, 3.65, 1.45, "模型范围", "从 MLP -> Tiny CNN -> QONNX -> 小型残差模块逐步推进。", "orange")
    add_card(slide, 4.85, 3.65, 3.65, 1.45, "论文/项目包装", "是否以 Agent 工程、EDA tool-use、HLS verification memory 作为创新点。", "blue")
    add_card(slide, 8.85, 3.65, 3.65, 1.45, "阶段验收", "每阶段保留真实 run、summary、trace、report、demo 视频。", "red")
    add_textbox(slide, 0.95, 6.15, 11.4, 0.38, "我希望导师支持：确定板卡/工具链方向，并认可先做“可验证模型到 HLS Agent”的工程主线。", size=16, bold=True, color=COLORS["navy"])
    add_footer(slide, 13)
    notes.append("这一页是说服导师的核心：你需要什么支持，以及导师投入后能得到什么阶段成果。")

    # 14
    slide = new_slide(prs, "下一步计划", "从可演示 MVP 走向更强研究成果")
    add_table(slide, 0.75, 1.35, 11.9, 4.7, 6, 4, [
        ["阶段", "目标", "可交付物", "验收指标"],
        ["1. 稳定 MNIST demo", "扩大验证样本", "100/1000 sample report", "accuracy/argmax/timing"],
        ["2. Precision Advisor", "降低 drift/resource", "per-layer precision 建议", "资源下降或误差下降"],
        ["3. CNN adapter", "接入 mnist-8/tiny CNN", "Conv/Pool/layout 支持", "真实 csim/csynth"],
        ["4. Benchmark", "Agent 能力量化", "workflow/RAG/repair 指标", "成功率和耗时"],
        ["5. 上板前准备", "AXI/IP packaging", "HLS IP + driver plan", "待板卡确认"],
    ], font_size=9)
    add_card(slide, 1.0, 6.15, 11.1, 0.75, "短期目标：把 MNIST 识别 demo 从 20-sample csim 扩展成完整 benchmark；中期目标：Tiny CNN 真实跑通；长期目标：上板验证。", "green")
    add_footer(slide, 14)
    notes.append("给导师清晰路线图：短中长期目标。")

    # 15
    slide = new_slide(prs, "结论：为什么现在值得继续投入", "已有真实结果，问题边界清楚，下一步有明确增量")
    add_card(slide, 0.9, 1.45, 3.65, 1.55, "已经证明可行", "真实 MNIST MLP 完成 hls4ml + Vivado HLS + recognition verification。", "green")
    add_card(slide, 4.85, 1.45, 3.65, 1.55, "具备 Agent 特色", "Todo/Specialist/ToolRegistry/Trace/Memory/RAG，不是单脚本。", "teal")
    add_card(slide, 8.8, 1.45, 3.65, 1.55, "有研究延展性", "CNN frontend、precision advisor、workflow benchmark、上板前 IP 化。", "orange")
    add_table(slide, 1.15, 3.85, 10.95, 1.45, 4, 3, [
        ["证明点", "当前证据", "下一步价值"],
        ["功能正确", "HLS acc 95%, argmax 100%", "扩大验证集"],
        ["综合可行", "timing met, report parsed", "资源优化"],
        ["Agent 工程", "trace/state/memory/specialists", "论文/实习项目亮点"],
    ], font_size=10)
    add_textbox(slide, 1.0, 6.25, 11.2, 0.45, "请求：希望导师支持我继续沿“模型到 HLS 的可验证 Agent”方向推进，并指导板卡/工具链/评价指标选择。", size=17, bold=True, color=COLORS["navy"], align=PP_ALIGN.CENTER)
    add_footer(slide, 15)
    notes.append("结尾直接给结论和请求。")

    prs.save(PPTX_PATH)

    notes_text = ["# HLS Agent 导师汇报 PPT 讲稿\n"]
    for idx, item in enumerate(notes, 1):
        notes_text.append(f"## Slide {idx}\n\n{item}\n")
    NOTES_PATH.write_text("\n".join(notes_text), encoding="utf-8")


if __name__ == "__main__":
    build_deck()
    print(PPTX_PATH)
    print(NOTES_PATH)
